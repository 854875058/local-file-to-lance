# -*- coding: utf-8 -*-
"""ETL：内容提取、管道处理、批量/SFTP 任务"""

import os
import uuid
import time
import logging
import shutil
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor, as_completed

import boto3
import pandas as pd
from PIL import Image
import docx
import pypdf
from pdf2image import convert_from_path
from pptx import Presentation
import paramiko

from config import (
    S3_CONFIG,
    TEMP_DIR,
    EXTRACT_DIR,
    CONTENT_EXTS,
    IMAGE_EXTS,
    ARCHIVE_EXTS,
    CHUNK_SIZE,
    CHUNK_OVERLAP,
    MAX_FILE_SIZE_MB,
)
from database import (
    calculate_file_hash,
    check_file_exists,
    register_file,
    insert_task_stat,
    delete_file_from_registry,
    insert_file_entities,
)
from models_loader import get_text_splitter

logger = logging.getLogger(__name__)

# S3 客户端（可选）
_s3_client = None

def _sanitize_filename(name: str) -> str:
    # S3 key 里避免出现路径分隔符
    return (name or "").replace("\\", "_").replace("/", "_")


def _category_for_ext(ext: str) -> str:
    e = (ext or "").lower()
    if e in IMAGE_EXTS:
        return "image"
    if e in ["mp3", "wav", "m4a", "flac", "ogg"]:
        return "audio"
    if e in ["mp4", "webm", "mov", "avi", "mkv"]:
        return "video"
    if e in ["pdf", "docx", "pptx", "txt", "md", "csv", "xlsx", "xls", "json", "log", "sql", "xml", "yaml", "ini", "py", "js", "sh"]:
        return "text"
    if e in ARCHIVE_EXTS:
        return "archive"
    return "other"


def get_s3_client():
    global _s3_client
    if _s3_client is None:
        try:
            _s3_client = boto3.client(
                "s3",
                endpoint_url=S3_CONFIG["endpoint_url"],
                aws_access_key_id=S3_CONFIG["access_key_id"],
                aws_secret_access_key=S3_CONFIG["secret_access_key"],
            )
            # 原始文件桶 & LanceDB 桶（LanceDB 桶主要给 lancedb 自己用，这里只确保存在）
            for b in [S3_CONFIG["raw_bucket"], S3_CONFIG["lance_bucket"]]:
                try:
                    _s3_client.create_bucket(Bucket=b)
                except Exception:
                    pass
        except Exception:
            _s3_client = False
    return _s3_client if _s3_client else None


def extract_content(path, ext, models):
    """全能内容提取：文本、文档、表格、音视频"""
    content = ""
    msg = ""

    try:
        if ext in ["mp3", "wav", "m4a", "mp4", "avi", "mov", "mkv", "flac"]:
            result = models["whisper"].transcribe(path)
            content = result.get("text", "")
            msg = "语音转录完成"

        elif ext in ["txt", "md", "py", "json", "log", "sh", "js", "java", "sql", "xml", "yaml", "ini"]:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()

        elif ext == "docx":
            content = "\n".join([p.text for p in docx.Document(path).paragraphs])
        elif ext == "pdf":
            content = "\n".join(
                [p.extract_text() or "" for p in pypdf.PdfReader(path).pages]
            )
        elif ext == "pptx":
            prs = Presentation(path)
            content = "\n".join(
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )
        elif ext == "csv":
            content = pd.read_csv(path).to_string()
        elif ext in ["xlsx", "xls"]:
            content = pd.read_excel(path).to_string()
        elif ext == "parquet":
            content = pd.read_parquet(path).to_string()
    except Exception as e:
        logger.error("提取失败 %s: %s", ext, e)
        msg = str(e)

    return content, msg


def process_pipeline(local_path, original_filename, models, tbl_text, tbl_image, tbl_files):
    if original_filename is None:
        original_filename = os.path.basename(local_path)
    ext = original_filename.rsplit(".", 1)[-1].lower() if "." in original_filename else ""

    overwrite = False
    try:
        f_hash = calculate_file_hash(local_path)
        overwrite = check_file_exists(f_hash)
        # SQLite 登记（重复则忽略）
        register_file(f_hash, original_filename, os.path.getsize(local_path))
    except Exception:
        # f_hash 仍可能用于后续流程；若异常，后面会再算一次
        pass

    # 压缩包
    if ext in ARCHIVE_EXTS:
        import zipfile
        import tarfile
        import shutil

        try:
            extract_folder = os.path.join(EXTRACT_DIR, str(uuid.uuid4()))
            os.makedirs(extract_folder)
            if ext == "zip":
                with zipfile.ZipFile(local_path, "r") as z:
                    z.extractall(extract_folder)
            else:
                with tarfile.open(local_path, "r") as t:
                    t.extractall(extract_folder)

            sub_files = []
            for root, _, files in os.walk(extract_folder):
                for f in files:
                    if not f.startswith("."):
                        sub_files.append((os.path.join(root, f), f))

            total = 0
            for p, n in sub_files:
                res = process_pipeline(p, n, models, tbl_text, tbl_image, tbl_files)
                if res["success"]:
                    total += res["count"]
            shutil.rmtree(extract_folder)
            return {"success": True, "msg": f"解压入库 {total} 文件", "count": total, "status": "ok"}
        except Exception as e:
            return {"success": False, "msg": str(e), "count": 0, "status": "error"}

    # 单文件
    try:
        f_hash = calculate_file_hash(local_path)
        if not f_hash:
            return {"success": False, "msg": "文件hash计算失败", "count": 0, "status": "error"}

        s3_client = get_s3_client()
        s3_uri = f"local://{original_filename}"
        if s3_client:
            try:
                safe_name = _sanitize_filename(original_filename)
                cat = _category_for_ext(ext)
                # S3 里"目录"本质是 key 前缀，按日期 + 类型分组更好管理
                today = datetime.now().strftime("%Y/%m/%d")
                key = f"raw/{today}/{cat}/{uuid.uuid4().hex[:8]}_{safe_name}"
                s3_client.upload_file(local_path, S3_CONFIG["raw_bucket"], key)
                s3_uri = f"s3://{S3_CONFIG['raw_bucket']}/{key}"
            except Exception as e:
                logger.warning(f"S3上传失败，使用本地URI: {e}")

        processed = False

        # 覆盖式重跑：同一 file_hash 先删旧记录，再写新记录（保证预览/检索一致）
        # 注意：为了避免删除后写入失败导致数据丢失，我们先准备好所有数据再删除
        if overwrite:
            logger.info(f"检测到重复文件，将覆盖: {original_filename}, hash={f_hash}")

        # 1) 原始文件入 LanceDB（用于整文件预览/下载）
        file_bytes = None
        try:
            # 检查文件大小，避免大文件占用过多内存
            file_size = os.path.getsize(local_path)
            max_file_size = MAX_FILE_SIZE_MB * 1024 * 1024  # 转换为字节

            if file_size > max_file_size:
                logger.warning(f"文件过大 ({file_size / 1024 / 1024:.2f}MB)，跳过存储到 files 表: {original_filename}")
                # 大文件只存储元数据，不存储 bytes
                file_row = {
                    "file_hash": f_hash,
                    "doc_name": original_filename,
                    "doc_type": ext,
                    "source_uri": s3_uri,
                    "file_bytes": b"",  # 空 bytes
                    "text_full": "",
                }
            else:
                with open(local_path, "rb") as rf:
                    file_bytes = rf.read()

                if not file_bytes:
                    logger.error(f"文件读取为空: {original_filename}")
                    return {"success": False, "msg": "文件读取为空", "count": 0, "status": "error"}

                file_row = {
                    "file_hash": f_hash,
                    "doc_name": original_filename,
                    "doc_type": ext,
                    "source_uri": s3_uri,
                    "file_bytes": file_bytes,
                    "text_full": "",
                }

            # 如果是覆盖模式，先删除旧的 files 表记录
            if overwrite:
                try:
                    # 转义单引号避免 SQL 注入
                    safe_hash = f_hash.replace("'", "''")
                    tbl_files.delete(f"file_hash = '{safe_hash}'")
                    logger.info(f"已删除旧的 files 表记录: hash={f_hash}")
                except Exception as e:
                    logger.warning(f"删除旧 files 表记录失败（可能不存在）: {e}")

            # 准备并写入 files 表数据
            tbl_files.add([file_row])
            logger.info(f"files 表写入成功: {original_filename}, hash={f_hash}, size={file_size} bytes")
        except Exception as e:
            logger.error(f"files 表写入失败: {e}, file={original_filename}, hash={f_hash}")
            import traceback
            logger.error(traceback.format_exc())
            # 如果 files 表写入失败，返回错误而不是继续处理
            return {"success": False, "msg": f"files表写入失败: {str(e)}", "count": 0, "status": "error"}

        # 2) 向量化入库（用于检索）
        if ext in CONTENT_EXTS:
            content, msg = extract_content(local_path, ext, models)
            if content and content.strip():
                # 如果是覆盖模式，先删除旧的 text 表记录
                if overwrite:
                    try:
                        # 转义单引号避免 SQL 注入
                        safe_hash = f_hash.replace("'", "''")
                        tbl_text.delete(f"file_hash = '{safe_hash}'")
                        logger.info(f"已删除旧的 text_chunks 表记录: hash={f_hash}")
                    except Exception as e:
                        logger.warning(f"删除旧 text_chunks 表记录失败: {e}")

                splitter = get_text_splitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)
                chunks = splitter.split_text(content)
                if chunks:
                    vecs = models["text"].encode(chunks)
                    data = [
                        {
                            "id": str(uuid.uuid4()),
                            "vector": v,
                            "text": c,
                            "source_uri": s3_uri,
                            "doc_name": original_filename,
                            "doc_type": ext,
                            "file_hash": f_hash,  # 直接写入，表一定有此列
                        }
                        for c, v in zip(chunks, vecs)
                    ]
                    tbl_text.add(data)
                    logger.info(f"text_chunks 表写入成功: {len(chunks)} 个切片, hash={f_hash}")
                    # 同步全文到 files 表（便于"整份文档"预览）
                    # 这里用 update（若版本不支持则忽略，仍可下载原件）
                    try:
                        safe_hash = f_hash.replace("'", "''")
                        tbl_files.update(where=f"file_hash = '{safe_hash}'", values={"text_full": content})
                        logger.info(f"files 表 text_full 更新成功: hash={f_hash}")
                    except Exception as e:
                        logger.warning(f"files 表 text_full 更新失败: {e}")
                    processed = True

        if ext in IMAGE_EXTS:
            try:
                # 如果是覆盖模式，先删除旧的 image 表记录
                if overwrite:
                    try:
                        safe_hash = f_hash.replace("'", "''")
                        tbl_image.delete(f"file_hash = '{safe_hash}'")
                        logger.info(f"已删除旧的 image_chunks 表记录: hash={f_hash}")
                    except Exception as e:
                        logger.warning(f"删除旧 image_chunks 表记录失败: {e}")

                img = Image.open(local_path)
                vec = models["clip_vision"].encode(img)
                row = {
                    "id": str(uuid.uuid4()),
                    "vector": vec,
                    "source_uri": s3_uri,
                    "doc_name": original_filename,
                    "meta_info": "image_file",
                    "file_hash": f_hash,  # 直接写入，表一定有此列
                }
                tbl_image.add([row])
                logger.info(f"image_chunks 表写入成功: {original_filename}, hash={f_hash}")
                processed = True
            except Exception as e:
                logger.warning(f"图像向量化失败: {e}")

        if ext == "pdf":
            try:
                # 如果是覆盖模式，先删除旧的 PDF 图像记录
                if overwrite:
                    try:
                        safe_hash = f_hash.replace("'", "''")
                        tbl_image.delete(f"file_hash = '{safe_hash}'")
                        logger.info(f"已删除旧的 PDF 图像记录: hash={f_hash}")
                    except Exception as e:
                        logger.warning(f"删除旧 PDF 图像记录失败: {e}")

                images = convert_from_path(local_path)
                if images:
                    vecs = models["clip_vision"].encode(images)
                    data = [
                        {
                            "id": str(uuid.uuid4()),
                            "vector": v,
                            "source_uri": s3_uri,
                            "doc_name": original_filename,
                            "meta_info": f"Page {i+1}",
                            "file_hash": f_hash,  # 直接写入，表一定有此列
                        }
                        for i, v in enumerate(vecs)
                    ]
                    tbl_image.add(data)
                    logger.info(f"PDF 图像向量化成功: {len(images)} 页, hash={f_hash}")
                    processed = True
            except Exception as e:
                logger.warning(f"PDF 图像向量化失败: {e}")

        if processed:
            # 方式B：不落本地预览目录，原始文件已写入 LanceDB `files` 表
            return {"success": True, "msg": ("覆盖OK" if overwrite else "OK"), "count": 1, "status": "ok"}
        return {"success": False, "msg": "Skipped", "count": 0, "status": "skipped"}
    except Exception as e:
        logger.exception("process_pipeline error: %s", e)
        return {"success": False, "msg": str(e), "count": 0, "status": "error"}


def delete_file_by_hash(file_hash, tbl_text, tbl_image, tbl_files):
    """删除文件的所有数据：file_registry + text_chunks + image_chunks + files 四张表"""
    safe_hash = file_hash.replace("'", "''")
    errors = []
    for tbl, name in [(tbl_text, 'text_chunks'), (tbl_image, 'image_chunks'), (tbl_files, 'files')]:
        try:
            tbl.delete(f"file_hash = '{safe_hash}'")
        except Exception as e:
            errors.append(f"{name}: {e}")
            logger.warning(f"删除 {name} 表记录失败: {e}")
    if not delete_file_from_registry(file_hash):
        errors.append("file_registry: 删除失败")
    if errors:
        logger.warning(f"删除文件 {file_hash} 部分失败: {errors}")
    return len(errors) == 0


def extract_entities_llm(text, file_hash):
    """调用 DeepSeek API 从文本中抽取实体，写入 file_entities 表。失败不阻塞主流程。"""
    try:
        from config import DEEPSEEK_API_KEY, DEEPSEEK_BASE_URL, DEEPSEEK_MODEL
        import json
        import urllib.request

        if not DEEPSEEK_API_KEY or not text or not text.strip():
            return

        # 截取前 3000 字符避免 token 过长
        snippet = text[:3000]
        prompt = (
            "请从以下文本中抽取关键实体（人名、地名、组织、技术术语等），"
            "以 JSON 数组格式返回，每个元素包含 name 和 type 两个字段。"
            "只返回 JSON 数组，不要其他内容。\n\n"
            f"文本：{snippet}"
        )

        payload = json.dumps({
            "model": DEEPSEEK_MODEL,
            "messages": [{"role": "user", "content": prompt}],
            "temperature": 0.1,
        }).encode("utf-8")

        req = urllib.request.Request(
            f"{DEEPSEEK_BASE_URL}/v1/chat/completions",
            data=payload,
            headers={
                "Content-Type": "application/json",
                "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
            },
        )
        with urllib.request.urlopen(req, timeout=30) as resp:
            result = json.loads(resp.read().decode("utf-8"))

        content = result["choices"][0]["message"]["content"].strip()
        # 尝试提取 JSON 数组
        if content.startswith("```"):
            content = content.split("```")[1]
            if content.startswith("json"):
                content = content[4:]
        entities_raw = json.loads(content)
        if not isinstance(entities_raw, list):
            return

        entities = []
        for e in entities_raw:
            name = (e.get("name") or "").strip()
            etype = (e.get("type") or "").strip()
            if name:
                entities.append((name, etype))

        if entities:
            insert_file_entities(file_hash, entities)
            logger.info(f"实体抽取完成: hash={file_hash}, 共 {len(entities)} 个实体")
    except Exception as e:
        logger.warning(f"实体抽取失败（不影响主流程）: {e}")


def batch_process_local_files(file_paths, models, tbl_text, tbl_image, tbl_files, progress_callback=None):
    """处理本地文件路径列表（NiceGUI 等非 Streamlit 前端使用）。
    file_paths: list of (local_path, original_filename) 元组
    返回: (succ, skip, dur, skipped_names)
    """
    start = time.time()
    total = len(file_paths)
    results = []
    skipped_names = []

    def process_one(item):
        local_path, name = item
        try:
            res = process_pipeline(local_path, name, models, tbl_text, tbl_image, tbl_files)
            # 异步实体抽取（成功入库的文本文件）
            if res.get("status") == "ok":
                try:
                    ext = name.split(".")[-1].lower()
                    if ext in CONTENT_EXTS:
                        content, _ = extract_content(local_path, ext, models)
                        if content and content.strip():
                            f_hash = calculate_file_hash(local_path)
                            extract_entities_llm(content, f_hash)
                except Exception:
                    pass
            return res, name
        except Exception as e:
            logger.error(f"处理文件失败: {name}, {e}")
            return {"success": False, "msg": str(e), "count": 0, "status": "error"}, name

    try:
        with ThreadPoolExecutor(max_workers=3) as executor:
            futures = {executor.submit(process_one, item): item for item in file_paths}
            for i, future in enumerate(as_completed(futures)):
                try:
                    res, name = future.result()
                    results.append(res)
                    if res.get("status") == "skipped":
                        skipped_names.append(name)
                    if progress_callback:
                        progress_callback(i + 1, total, res["msg"])
                except Exception as e:
                    logger.error(f"获取任务结果失败: {e}")
                    results.append({"success": False, "msg": str(e), "count": 0, "status": "error"})
    finally:
        pass  # 本地路径由调用方管理清理

    succ = sum(r["count"] for r in results if r["status"] == "ok")
    skip = sum(1 for r in results if r["status"] == "skipped")
    dur = time.time() - start
    insert_task_stat("batch", total, succ, dur)
    return succ, skip, dur, skipped_names


def sftp_task(host, port, user, password, path, models, tbl_text, tbl_image, tbl_files, progress_callback=None):
    logs = []
    skipped_names = []
    tr = None
    try:
        tr = paramiko.Transport((host, int(port)))
        tr.connect(username=user, password=password)
        sftp = paramiko.SFTPClient.from_transport(tr)
        files = sftp.listdir(path)
        logs.append(f"🔗 扫描到 {len(files)} 个文件")

        local_fs = []
        total_files = 0
        for f in files:
            if f.startswith("."):
                continue
            total_files += 1
            try:
                local_path = os.path.join(TEMP_DIR, f)
                sftp.get(path.rstrip("/") + "/" + f, local_path)
                local_fs.append((local_path, f))
                if progress_callback:
                    progress_callback(len(local_fs), total_files, f"下载: {f}")
            except Exception:
                logs.append(f"⚠️ 下载失败 {f}")

        cnt = 0
        for i, (local_path, name) in enumerate(local_fs):
            res = process_pipeline(local_path, name, models, tbl_text, tbl_image, tbl_files)
            if res["status"] == "ok":
                cnt += res["count"]
                # 异步实体抽取
                try:
                    ext = name.split(".")[-1].lower()
                    if ext in CONTENT_EXTS:
                        content, _ = extract_content(local_path, ext, models)
                        if content and content.strip():
                            f_hash = calculate_file_hash(local_path)
                            extract_entities_llm(content, f_hash)
                except Exception:
                    pass
            elif res["status"] == "skipped":
                skipped_names.append(name)
            if progress_callback:
                progress_callback(i + 1, len(local_fs), f"处理: {name}")
            if os.path.exists(local_path):
                os.remove(local_path)
        logs.append(f"🎉 入库 {cnt} 条")
        if skipped_names:
            logs.append(f"⏭️ 跳过 {len(skipped_names)} 个文件: {', '.join(skipped_names)}")
    except Exception as e:
        logs.append(f"🔴 {e}")
    finally:
        if tr:
            tr.close()
    return logs, skipped_names
